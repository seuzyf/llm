import os
import tempfile
import zipfile
import shutil
from .utils import indent

def _progid_to_ext(prog_id: str) -> str:
    prog_id = (prog_id or '').lower()
    mapping = {
        'word.document': '.docx', 'word.sheet': '.docx',
        'excel.sheet': '.xlsx', 'excel.chart': '.xlsx',
        'powerpoint': '.pptx', 'acrord': '.pdf',
        'pdf': '.pdf', 'package': '',
    }
    for key, ext in mapping.items():
        if key in prog_id:
            return ext
    return ''

def _extract_pptx_slide_attachments(slide, pptx_zip, slide_idx):
    attachments = []
    slide_part_name = slide.part.partname
    rel_path = slide_part_name.lstrip('/')
    rel_dir  = os.path.dirname(rel_path)
    rels_path = f"ppt/slides/_rels/slide{slide_idx}.xml.rels"
    
    try:
        with pptx_zip.open(rels_path) as f:
            import xml.etree.ElementTree as ET
            tree = ET.parse(f)
            root = tree.getroot()
            ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
            
            for rel in root.findall('r:Relationship', ns):
                rel_type = rel.get('Type', '')
                target   = rel.get('Target', '')
                if 'oleObject' not in rel_type and 'package' not in rel_type:
                    continue
                if target.startswith('/'):
                    zip_inner_path = target.lstrip('/')
                else:
                    zip_inner_path = os.path.normpath(os.path.join(rel_dir, target)).replace('\\', '/')
                
                if zip_inner_path in pptx_zip.namelist():
                    data = pptx_zip.read(zip_inner_path)
                    display_name = os.path.basename(zip_inner_path)
                    attachments.append((display_name, data))
    except (KeyError, Exception):
        pass

    try:
        import xml.etree.ElementTree as ET
        named_map = {}
        slide_xml_path = rel_path
        with pptx_zip.open(slide_xml_path) as f:
            slide_tree = ET.parse(f)
            
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        for ole in slide_tree.iter('{http://schemas.openxmlformats.org/presentationml/2006/main}oleObj'):
            r_id = ole.get(f'{{{ns_r}}}id') or ole.get('r:id') or ''
            prog_id = ole.get('progId') or ''
            user_label = ole.get('name') or ''
            ext_guess = _progid_to_ext(prog_id)
            if ext_guess and user_label:
                named_map[r_id] = f"{user_label}{ext_guess}"
        
        with pptx_zip.open(rels_path) as f:
            tree = ET.parse(f)
            root = tree.getroot()
            ns_rel = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
            id_to_zip = {}
            for rel in root.findall('r:Relationship', ns_rel):
                rel_type = rel.get('Type', '')
                if 'oleObject' in rel_type or 'package' in rel_type:
                    target = rel.get('Target', '')
                    zpath = target.lstrip('/') if target.startswith('/') else os.path.normpath(os.path.join(rel_dir, target)).replace('\\', '/')
                    id_to_zip[rel.get('Id', '')] = zpath
            
            new_attachments = []
            for (dname, data) in attachments:
                replaced = False
                for rid, friendly in named_map.items():
                    if id_to_zip.get(rid) and os.path.basename(id_to_zip[rid]) == dname:
                        new_attachments.append((friendly, data))
                        replaced = True
                        break
                if not replaced:
                    new_attachments.append((dname, data))
            attachments = new_attachments
    except Exception:
        pass

    return attachments

def parse_ppt(file_path, display_name, parse_callback):
    ext = os.path.splitext(display_name)[1].lower()
    
    if ext == '.pptx':
        from pptx import Presentation
        prs = Presentation(file_path)
        slide_outputs = []

        with zipfile.ZipFile(file_path, 'r') as pptx_zip:
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_parts = [f"── 第 {slide_idx} 页 ──"]

                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
                if text_runs:
                    slide_parts.append("\n".join(text_runs))

                try:
                    embedded = _extract_pptx_slide_attachments(slide, pptx_zip, slide_idx)
                except Exception:
                    embedded = []

                if embedded:
                    slide_parts.append(f"\n  📦 本页包含 {len(embedded)} 个附件：")
                    for att_name, att_data in embedded:
                        att_ext = os.path.splitext(att_name)[1] or '.bin'
                        fd, tmp_path = tempfile.mkstemp(suffix=att_ext)
                        os.close(fd)
                        try:
                            with open(tmp_path, 'wb') as f:
                                f.write(att_data)
                            att_text = parse_callback(tmp_path, display_name=att_name)
                            slide_parts.append(f"\n  📎 附件 [{att_name}]:\n" + indent(att_text, prefix="    "))
                        except Exception as e:
                            slide_parts.append(f"\n  ❌ 附件 [{att_name}] 解析跳过: {e}")
                        finally:
                            if os.path.exists(tmp_path):
                                os.remove(tmp_path)

                slide_outputs.append("\n".join(slide_parts))

        return "\n\n".join(slide_outputs)

    elif ext == '.ppt':
        try:
            import win32com.client
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            abs_path = os.path.abspath(file_path)
            prs_com = powerpoint.Presentations.Open(abs_path, WithWindow=False)

            tmp_dir = tempfile.mkdtemp()
            tmp_pptx = os.path.join(tmp_dir, "converted.pptx")
            prs_com.SaveAs(tmp_pptx, FileFormat=24)
            prs_com.Close()
            powerpoint.Quit()

            return parse_callback(tmp_pptx, display_name=display_name)
        except Exception as e:
            raise RuntimeError(f"读取 .ppt 失败: {str(e)}")
        finally:
            try:
                shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception:
                pass