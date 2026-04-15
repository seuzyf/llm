import sys
import os
import io
import email
import tempfile
import zipfile
import shutil

try:
    import extract_msg
except ImportError:
    extract_msg = None

# 安全地设置 UTF-8 编码，避免重复包裹导致底层流被关闭
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
if hasattr(sys.stderr, 'reconfigure'):
    sys.stderr.reconfigure(encoding='utf-8', errors='replace')


# ──────────────────────────────────────────────
# 工具函数：从 pptx Shape 中提取嵌入附件
# ──────────────────────────────────────────────
def _extract_pptx_slide_attachments(slide, pptx_zip, slide_idx):
    """
    扫描单张幻灯片的所有 Shape，将 OLE 对象 / 打包对象 / 超链接文件 提取出来。
    返回列表：[(display_name, bytes_data), ...]
    
    pptx 本质是 ZIP，内部结构：
      ppt/slides/slide{n}.xml          —— 幻灯片主体
      ppt/slides/_rels/slide{n}.xml.rels —— 关系文件（含附件引用）
      ppt/embeddings/                  —— 嵌入的 OLE/文件对象
    """
    attachments = []
    
    # 关系文件路径
    slide_part_name = slide.part.partname          # e.g. /ppt/slides/slide1.xml
    # python-pptx partname 带前导 /，zip 内无前导 /
    rel_path = slide_part_name.lstrip('/')         # ppt/slides/slide1.xml
    rel_dir  = os.path.dirname(rel_path)           # ppt/slides
    rels_path = f"ppt/slides/_rels/slide{slide_idx}.xml.rels"
    
    # ── 方法①：从 ZIP 关系文件里找 oleObject / package 类型的附件 ──
    try:
        with pptx_zip.open(rels_path) as f:
            import xml.etree.ElementTree as ET
            tree = ET.parse(f)
            root = tree.getroot()
            ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
            
            for rel in root.findall('r:Relationship', ns):
                rel_type = rel.get('Type', '')
                target   = rel.get('Target', '')
                
                # 只关心嵌入对象（oleObject、package、image 不在此列）
                if 'oleObject' not in rel_type and 'package' not in rel_type:
                    continue
                
                # target 可能是相对路径，如 ../embeddings/oleObject1.bin
                if target.startswith('/'):
                    zip_inner_path = target.lstrip('/')
                else:
                    zip_inner_path = os.path.normpath(
                        os.path.join(rel_dir, target)
                    ).replace('\\', '/')
                
                if zip_inner_path in pptx_zip.namelist():
                    data = pptx_zip.read(zip_inner_path)
                    display_name = os.path.basename(zip_inner_path)
                    attachments.append((display_name, data))
    except KeyError:
        pass  # 该幻灯片无 rels 文件，正常跳过
    except Exception:
        pass

    # ── 方法②：从 Shape 的 XML 属性里读取原始文件名（改善展示名称）──
    # python-pptx 的 shape.element 可拿到原始 XML，从中找 progId / userLabel
    try:
        import xml.etree.ElementTree as ET
        named_map = {}   # zip_inner_path -> friendly_name
        
        slide_xml_path = rel_path
        with pptx_zip.open(slide_xml_path) as f:
            slide_tree = ET.parse(f)
        
        # 收集所有 oleObj 节点的 r:id -> progId/name 映射
        ns_p   = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns_a   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        ns_v   = 'urn:schemas-microsoft-com:vml'
        
        for ole in slide_tree.iter('{http://schemas.openxmlformats.org/presentationml/2006/main}oleObj'):
            r_id      = ole.get(f'{{{ns_r}}}id') or ole.get('r:id') or ''
            prog_id   = ole.get('progId') or ''
            user_label = ole.get('name') or ''
            
            # 根据 progId 猜扩展名，使展示名更友好
            ext_guess = _progid_to_ext(prog_id)
            
            # 用 r_id 在 rels 里查真实 zip 路径（已在方法①里拿到了 data）
            # 这里只更新 display_name，与方法①采集的顺序对应
            if ext_guess and user_label:
                key_name = f"{user_label}{ext_guess}"
                named_map[r_id] = key_name
        
        # 将友好名字回写到 attachments（按 r:id 对应关系重新建立）
        # 重新扫一次 rels，建立 r:id -> zip_path 映射
        with pptx_zip.open(rels_path) as f:
            tree = ET.parse(f)
            root = tree.getroot()
            ns_rel = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
            
            id_to_zip = {}
            for rel in root.findall('r:Relationship', ns_rel):
                rel_type = rel.get('Type', '')
                target   = rel.get('Target', '')
                rid      = rel.get('Id', '')
                if 'oleObject' not in rel_type and 'package' not in rel_type:
                    continue
                if target.startswith('/'):
                    zpath = target.lstrip('/')
                else:
                    zpath = os.path.normpath(
                        os.path.join(rel_dir, target)
                    ).replace('\\', '/')
                id_to_zip[rid] = zpath
            
            # 用友好名替换 attachments 里的 display_name
            new_attachments = []
            for (dname, data) in attachments:
                replaced = False
                for rid, friendly in named_map.items():
                    if id_to_zip.get(rid) and os.path.basename(id_to_zip[rid]) == dname:
                        new_attachments.append((friendly, data))
                        replaced = True
                        break
                if not replaced:
                    new_attachments.append((dname, data))
            attachments = new_attachments
            
    except Exception:
        pass  # 友好名替换失败，保留原始名称即可

    return attachments


def _progid_to_ext(prog_id: str) -> str:
    """根据 OLE progId 猜测文件扩展名"""
    prog_id = (prog_id or '').lower()
    mapping = {
        'word.document'    : '.docx',
        'word.sheet'       : '.docx',
        'excel.sheet'      : '.xlsx',
        'excel.chart'      : '.xlsx',
        'powerpoint'       : '.pptx',
        'acrord'           : '.pdf',
        'pdf'              : '.pdf',
        'package'          : '',      # 通用包，不加扩展名
    }
    for key, ext in mapping.items():
        if key in prog_id:
            return ext
    return ''


# ──────────────────────────────────────────────
# 主解析函数
# ──────────────────────────────────────────────
def parse_file_to_text(file_path, display_name=None):
    """通用解析管道：返回提取出的纯文本，支持递归调用"""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    actual_name = display_name or os.path.basename(file_path)
    ext = os.path.splitext(actual_name)[1].lower()

    # ── 1. 邮件解析 (正文 + 附件递归) ──────────────────────────────────
    if ext in ['.eml', '.msg']:
        output = [f"📧 邮件: {actual_name}"]
        attachments = []

        if ext == '.eml':
            with open(file_path, 'rb') as f:
                msg = email.message_from_binary_file(f)

            body = ""
            for part in msg.walk():
                content_type = part.get_content_type()
                disp = str(part.get('Content-Disposition'))

                if content_type in ["text/plain", "text/html"] and "attachment" not in disp:
                    try:
                        payload = part.get_payload(decode=True).decode(
                            part.get_content_charset() or 'utf-8', errors='ignore'
                        )
                        if content_type == "text/html":
                            from bs4 import BeautifulSoup
                            payload = BeautifulSoup(
                                payload, 'html.parser'
                            ).get_text(separator='\n', strip=True)
                        body += payload + "\n"
                    except:
                        pass

            if body.strip():
                output.extend(["\n--- 邮件正文 ---", body.strip(), "-----------------\n"])

            for part in msg.walk():
                if (part.get_content_maintype() == 'multipart'
                        or part.get('Content-Disposition') is None):
                    continue
                filename = part.get_filename()
                if filename:
                    attachments.append((filename, part.get_payload(decode=True)))

        elif ext == '.msg':
            if extract_msg is None:
                return "❌ 缺少 extract-msg 库，无法解析 .msg"
            msg = extract_msg.Message(file_path)
            if msg.body:
                output.extend(["\n--- 邮件正文 ---", msg.body, "-----------------\n"])
            for attachment in msg.attachments:
                filename = attachment.longFilename or attachment.shortFilename
                if filename:
                    attachments.append((filename, attachment.data))

        # 递归处理所有附件
        if attachments:
            output.append(f"📦 发现 {len(attachments)} 个附件，开始自动提取...")
            for filename, payload in attachments:
                fd, tmp_path = tempfile.mkstemp(suffix=os.path.splitext(filename)[1])
                os.close(fd)
                try:
                    with open(tmp_path, 'wb') as f:
                        f.write(payload)
                    res = parse_file_to_text(tmp_path, display_name=filename)
                    if res:
                        output.append(f"\n📎 附件 [{filename}]:\n{res}")
                except Exception as e:
                    output.append(f"❌ 附件 [{filename}] 解析跳过: {str(e)}")
                finally:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)

        return "\n".join(output)

    # ── 2. 二进制文件拦截 ────────────────────────────────────────────────
    binary_exts = [
        '.exe', '.dll', '.so', '.bin',
        '.zip', '.rar', '.7z', '.tar', '.gz',
        '.jpg', '.jpeg', '.png', '.gif',
        '.mp3', '.mp4',
    ]
    if ext in binary_exts:
        raise ValueError(f"暂不支持解析媒体或二进制文件: {ext}")

    # ── 3. PDF 解析 ──────────────────────────────────────────────────────
    if ext == '.pdf':
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        return "\n".join(
            [p.extract_text() for p in reader.pages if p.extract_text()]
        )

    # ── 4. Word 解析 ─────────────────────────────────────────────────────
    elif ext == '.docx':
        from docx import Document
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    elif ext == '.doc':
        try:
            import win32com.client
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            abs_path = os.path.abspath(file_path)
            doc = word.Documents.Open(abs_path)
            text = doc.Content.Text
            doc.Close()
            word.Quit()
            return text
        except Exception as e:
            raise RuntimeError(f"读取 .doc 失败: {str(e)}")

    # ── 5. PPT 解析（含附件提取）────────────────────────────────────────
    elif ext == '.pptx':
        from pptx import Presentation
        prs = Presentation(file_path)
        slide_outputs = []

        # pptx 本质是 ZIP，同时打开备用（用于提取嵌入附件）
        with zipfile.ZipFile(file_path, 'r') as pptx_zip:
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_parts = [f"── 第 {slide_idx} 页 ──"]

                # ① 提取幻灯片文字
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
                if text_runs:
                    slide_parts.append("\n".join(text_runs))

                # ② 提取该页嵌入附件并递归解析
                try:
                    embedded = _extract_pptx_slide_attachments(
                        slide, pptx_zip, slide_idx
                    )
                except Exception:
                    embedded = []

                if embedded:
                    slide_parts.append(
                        f"\n  📦 本页包含 {len(embedded)} 个附件："
                    )
                    for att_name, att_data in embedded:
                        att_ext = os.path.splitext(att_name)[1] or '.bin'
                        fd, tmp_path = tempfile.mkstemp(suffix=att_ext)
                        os.close(fd)
                        try:
                            with open(tmp_path, 'wb') as f:
                                f.write(att_data)
                            att_text = parse_file_to_text(
                                tmp_path, display_name=att_name
                            )
                            slide_parts.append(
                                f"\n  📎 附件 [{att_name}]:\n"
                                + _indent(att_text, prefix="    ")
                            )
                        except Exception as e:
                            slide_parts.append(
                                f"\n  ❌ 附件 [{att_name}] 解析跳过: {e}"
                            )
                        finally:
                            if os.path.exists(tmp_path):
                                os.remove(tmp_path)

                slide_outputs.append("\n".join(slide_parts))

        return "\n\n".join(slide_outputs)

    elif ext == '.ppt':
        """
        .ppt 为旧版二进制格式，通过 win32com 将其另存为 .pptx
        再交给上方 .pptx 流程处理，保证附件逻辑一致。
        """
        try:
            import win32com.client
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            abs_path   = os.path.abspath(file_path)
            prs_com    = powerpoint.Presentations.Open(abs_path, WithWindow=False)

            # 另存为 pptx 到临时目录
            tmp_dir   = tempfile.mkdtemp()
            tmp_pptx  = os.path.join(tmp_dir, "converted.pptx")
            # FileFormat=24 => ppSaveAsOpenXMLPresentation (.pptx)
            prs_com.SaveAs(tmp_pptx, FileFormat=24)
            prs_com.Close()
            powerpoint.Quit()

            # 交给 .pptx 流程（含附件提取）
            result = parse_file_to_text(tmp_pptx, display_name=actual_name)
            return result
        except Exception as e:
            raise RuntimeError(f"读取 .ppt 失败: {str(e)}")
        finally:
            # 清理临时目录
            try:
                shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception:
                pass

    # ── 6. 表格解析 ──────────────────────────────────────────────────────
    elif ext in ['.xlsx', '.xls', '.csv']:
        import pandas as pd
        df = None
        if ext == '.csv':
            for enc in ['utf-8', 'gbk', 'utf-16', 'gb2312', 'latin1']:
                try:
                    df = pd.read_csv(file_path, encoding=enc)
                    break
                except:
                    continue
        else:
            try:
                if ext == '.xlsx':
                    df = pd.read_excel(file_path, engine="openpyxl")
                else:
                    df = pd.read_excel(file_path, engine="xlrd")
            except Exception:
                df = pd.read_excel(file_path)

        df_clean  = df.fillna("").astype(str)
        md_lines  = [f"### 📄 表格：{actual_name}", ""]

        if df_clean.empty:
            md_lines.append("（空表格）")
        else:
            header    = "| " + " | ".join(
                [str(col).replace('\n', ' ') for col in df_clean.columns]
            ) + " |"
            separator = "| " + " | ".join(
                ["---"] * len(df_clean.columns)
            ) + " |"
            md_lines.append(header)
            md_lines.append(separator)
            for _, row in df_clean.iterrows():
                row_str = "| " + " | ".join(
                    [str(cell).replace('\n', '<br>') for cell in row]
                ) + " |"
                md_lines.append(row_str)

        return "\n".join(md_lines)

    # ── 7. 网页及 XML ────────────────────────────────────────────────────
    elif ext in ['.html', '.htm', '.xml']:
        from bs4 import BeautifulSoup
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            return soup.get_text(separator='\n', strip=True)

    # ── 8. 富文本 ────────────────────────────────────────────────────────
    elif ext == '.rtf':
        from striprtf.striprtf import rtf_to_text
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return rtf_to_text(f.read())

    # ── 9. 电子书 ────────────────────────────────────────────────────────
    elif ext == '.epub':
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book  = epub.read_epub(file_path)
        texts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_body_content(), 'html.parser')
            texts.append(soup.get_text(separator='\n', strip=True))
        return "\n".join(texts)

    # ── 10. 兜底纯文本 ──────────────────────────────────────────────────
    else:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='gbk', errors='ignore') as f:
                return f.read()


# ──────────────────────────────────────────────
# 辅助：给多行文本统一缩进
# ──────────────────────────────────────────────
def _indent(text: str, prefix: str = "  ") -> str:
    return "\n".join(prefix + line for line in text.splitlines())


# ──────────────────────────────────────────────
# 入口
# ──────────────────────────────────────────────
def parse():
    if len(sys.argv) < 2:
        return

    file_path = sys.argv[1]
    try:
        result = parse_file_to_text(file_path)
        if result:
            print(result)
    except Exception as e:
        sys.stderr.write(f"解析异常: {str(e)}")
        sys.exit(1)


if __name__ == "__main__":
    parse()
